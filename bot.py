import asyncio, os, re, json, uuid, base64, time, logging, sys, random
from io import BytesIO
from datetime import datetime
from typing import Optional, Dict, List
from pathlib import Path

import aiohttp
from aiogram import Bot, Dispatcher, F
from aiogram.filters import Command, StateFilter
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.types import (
    BufferedInputFile, ReplyKeyboardMarkup, KeyboardButton,
    InlineKeyboardMarkup, InlineKeyboardButton, CallbackQuery, Message
)
from aiogram.client.session.aiohttp import AiohttpSession
from aiogram.exceptions import TelegramRetryAfter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from yookassa import Configuration, Payment

# ===== ЛОГИ =====
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
log = logging.getLogger("bot")

# ===== ЗАГРУЗКА .env =====
if not os.getenv("RAILWAY_ENVIRONMENT"):
    env_file = Path(__file__).parent / '.env'
    if env_file.exists():
        with open(env_file) as f:
            for line in f:
                line = line.strip()
                if line and not line.startswith('#') and '=' in line:
                    k, v = line.split('=', 1)
                    os.environ[k.strip()] = v.strip()

# ===== КЛЮЧИ =====
BOT_TOKEN = os.getenv("BOT_TOKEN")
GIGA_AUTH = base64.b64encode(
    f"{os.getenv('GIGACHAT_CLIENT_ID')}:{os.getenv('GIGACHAT_SECRET')}".encode()
).decode() if os.getenv('GIGACHAT_CLIENT_ID') else ""
YOOKASSA_ID = os.getenv("YOOKASSA_SHOP_ID")
YOOKASSA_KEY = os.getenv("YOOKASSA_SECRET_KEY")
UNSPLASH_KEY = os.getenv("UNSPLASH_ACCESS_KEY", "")
PIXABAY_KEY = os.getenv("PIXABAY_API_KEY", "")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "")
GOOGLE_CX = os.getenv("GOOGLE_CX", "")
ADMIN_ID = int(os.getenv("ADMIN_ID", "0"))
PRICE = 100

if not all([BOT_TOKEN, YOOKASSA_ID, YOOKASSA_KEY]):
    raise SystemExit("❌ Нет обязательных ключей!")

Configuration.account_id = YOOKASSA_ID
Configuration.secret_key = YOOKASSA_KEY

# ===== РАЗМЕРЫ СЛАЙДА =====
SW = Emu(12192000)
SH = Emu(6858000)
MG = Emu(365760)
GP = Emu(274320)
IW = Emu(4937760)
TW = SW - IW - GP - MG*2

# ===== СТИЛИ POWERPOINT =====
PPT_THEMES = [
    {"name": "Classic Blue", "bg": "1A3C6D", "accent": "F4A261", "text": "FFFFFF", "title": "FFFFFF"},
    {"name": "Dark Elegance", "bg": "2C3E50", "accent": "E74C3C", "text": "ECF0F1", "title": "E74C3C"},
    {"name": "Forest Green", "bg": "1B4332", "accent": "FFB703", "text": "E9F5DB", "title": "FFB703"},
    {"name": "Burgundy Gold", "bg": "3D0C11", "accent": "D4AF37", "text": "F5E6CC", "title": "D4AF37"},
    {"name": "Ocean Depth", "bg": "023047", "accent": "FB8500", "text": "8ECAE6", "title": "FB8500"},
    {"name": "Plum Velvet", "bg": "3C1642", "accent": "F72585", "text": "E5D9F2", "title": "F72585"},
    {"name": "Slate Modern", "bg": "334155", "accent": "38BDF8", "text": "F1F5F9", "title": "38BDF8"},
    {"name": "Espresso Cream", "bg": "3C2A21", "accent": "D4A574", "text": "F5E6D3", "title": "D4A574"},
]

# ===== СОСТОЯНИЯ =====
class State(StatesGroup):
    topic = State()
    payment = State()

bot = Bot(token=BOT_TOKEN, session=AiohttpSession(timeout=120))
dp = Dispatcher(storage=MemoryStorage())

def menu():
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text="🎨 Создать презентацию")],
        [KeyboardButton(text="ℹ️ Помощь"), KeyboardButton(text="💰 Цена")]
    ], resize_keyboard=True)

def pay_kb(url):
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💳 Оплатить 100₽", url=url)],
        [InlineKeyboardButton(text="✅ Я оплатил", callback_data="paid")],
        [InlineKeyboardButton(text="❌ Отмена", callback_data="cancel")]
    ])

# ===== HEX → RGB =====
def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

# ===== ПАРСИНГ JSON =====
def extract_json(text):
    if not text:
        return None
    start = text.find('{')
    end = text.rfind('}')
    if start == -1 or end == -1 or end <= start:
        return None
    json_str = text[start:end+1]
    in_string = False
    cleaned = []
    for c in json_str:
        if c == '"' and (not cleaned or cleaned[-1] != '\\'):
            in_string = not in_string
        if c == '\n' and in_string:
            cleaned.append('\\n')
        elif c == '\r' and in_string:
            continue
        elif c == '\t' and in_string:
            cleaned.append(' ')
        else:
            cleaned.append(c)
    json_str = ''.join(cleaned)
    json_str = re.sub(r',\s*}', '}', json_str)
    json_str = re.sub(r',\s*]', ']', json_str)
    try:
        return json.loads(json_str)
    except:
        return None

# ===== ТОКЕН GIGACHAT =====
_token = None
_token_exp = 0
_token_lock = asyncio.Lock()

async def get_giga_token():
    global _token, _token_exp
    if not GIGA_AUTH:
        return None
    async with _token_lock:
        if _token and time.time() < _token_exp - 300:
            return _token
        try:
            connector = aiohttp.TCPConnector(ssl=False)
            async with aiohttp.ClientSession(connector=connector) as s:
                async with s.post(
                    "https://ngw.devices.sberbank.ru:9443/api/v2/oauth",
                    headers={
                        "Authorization": f"Basic {GIGA_AUTH}",
                        "RqUID": str(uuid.uuid4()),
                        "Content-Type": "application/x-www-form-urlencoded"
                    },
                    data={"scope": "GIGACHAT_API_PERS"},
                    timeout=aiohttp.ClientTimeout(total=15)
                ) as r:
                    if r.status == 200:
                        data = await r.json()
                        _token = data["access_token"]
                        _token_exp = time.time() + 3600
                        return _token
        except:
            pass
        return None

async def ask_gigachat(prompt):
    if not GIGA_AUTH:
        return None
    token = await get_giga_token()
    if not token:
        return None
    try:
        connector = aiohttp.TCPConnector(ssl=False)
        async with aiohttp.ClientSession(connector=connector) as s:
            async with s.post(
                "https://gigachat.devices.sberbank.ru/api/v1/chat/completions",
                headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
                json={
                    "model": "GigaChat",
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": 0.7, "max_tokens": 4000
                },
                timeout=aiohttp.ClientTimeout(total=90)
            ) as r:
                if r.status == 200:
                    return (await r.json())["choices"][0]["message"]["content"]
    except:
        pass
    return None

async def ask_duckduckgo(prompt):
    try:
        async with aiohttp.ClientSession() as s:
            async with s.get("https://duckduckgo.com/duckchat/v1/status",
                           headers={"x-vqd-accept": "1"}, timeout=10) as r:
                if r.status != 200:
                    return None
                vqd = r.headers.get("x-vqd-4")
            async with s.post("https://duckduckgo.com/duckchat/v1/chat",
                            headers={"x-vqd-4": vqd, "Content-Type": "application/json"},
                            json={"model": "gpt-4o-mini", "messages": [{"role": "user", "content": prompt}]},
                            timeout=90) as r:
                if r.status == 200:
                    text = await r.text()
                    message = ""
                    for line in text.split('\n'):
                        if line.startswith('data: '):
                            try:
                                d = json.loads(line[6:])
                                if d.get("message"):
                                    message += d["message"]
                            except:
                                pass
                    if message:
                        return message.strip()
    except:
        pass
    return None

async def ask_ai(prompt):
    resp = await ask_gigachat(prompt)
    if resp: return resp
    resp = await ask_duckduckgo(prompt)
    if resp: return resp
    return None

# ===== ГЕНЕРАЦИЯ КОНТЕНТА =====
async def get_content(topic, user_n=None):
    if user_n:
        if user_n <= 5:
            depth = f"Ровно {user_n} слайдов. Только ключевые тезисы."
        elif user_n <= 8:
            depth = f"Ровно {user_n} слайдов. Факты + примеры."
        else:
            depth = f"Ровно {user_n} слайдов. Глубокое погружение."
    else:
        depth = "Определи оптимальное количество (5-8)."
    
    prompt = f"""Ты — Ведущий Архитектор Презентаций. Создай презентацию на РУССКОМ ЯЗЫКЕ на тему "{topic}".

{depth}

СТРУКТУРА:
- Слайды с 1 по предпоследний: содержательные. Каждый — ОДИН инсайт.
- Для каждого слайда:
  · "title": Заголовок-тезис НА РУССКОМ (4-8 слов)
  · "expert_text": 3-4 предложения НА РУССКОМ. Сценарий речи. Факты. Цифры.
  · "visual_prompt": Визуальная метафора НА АНГЛИЙСКОМ. Стиль + объект + свет + текстура.
    ВСЕГДА заканчивай на: «--no text, words, letters, logos, watermarks --ar 16:9»

ПРИМЕРЫ visual_prompt:
- «macro photography, ancient stone with fresh sprout, dramatic side light, dark background --no text --ar 16:9»
- «aerial view, lone tree on cliff at sunrise, volumetric fog, golden ratio --no text --ar 16:9»

ФОРМАТ ОТВЕТА — ТОЛЬКО JSON:
{{
  "topic": "Тема",
  "about_text": "2-3 предложения: О ЧЁМ эта презентация. Самая суть.",
  "slides": [
    {{"title": "...", "expert_text": "...", "visual_prompt": "..."}}
  ]
}}"""
    
    resp = await ask_ai(prompt)
    if not resp:
        return None
    data = extract_json(resp)
    if not data:
        resp2 = await ask_ai("Return ONLY JSON. " + prompt)
        if resp2:
            data = extract_json(resp2)
    return data if data and "slides" in data else None

# ===== ПОИСК КАРТИНОК =====

async def _download(url: str) -> Optional[BytesIO]:
    try:
        async with aiohttp.ClientSession() as s:
            async with s.get(url, timeout=20) as r:
                if r.status == 200:
                    data = await r.read()
                    if len(data) > 2000:
                        return BytesIO(data)
    except:
        pass
    return None

async def search_unsplash(query: str) -> Optional[BytesIO]:
    if not UNSPLASH_KEY or not query: return None
    try:
        async with aiohttp.ClientSession() as s:
            async with s.get(
                "https://api.unsplash.com/search/photos",
                params={"query": query, "per_page": 5, "orientation": "landscape", "client_id": UNSPLASH_KEY},
                timeout=15
            ) as r:
                if r.status == 200:
                    results = (await r.json()).get("results", [])
                    if results:
                        url = random.choice(results)["urls"].get("regular")
                        if url:
                            img = await _download(url)
                            if img:
                                log.info(f"✅ Unsplash: {query[:60]}")
                                return img
    except: pass
    return None

async def search_pixabay(query: str) -> Optional[BytesIO]:
    if not PIXABAY_KEY or not query: return None
    try:
        async with aiohttp.ClientSession() as s:
            for img_type in ["photo", "illustration"]:
                async with s.get(
                    "https://pixabay.com/api/",
                    params={"key": PIXABAY_KEY, "q": query, "per_page": 5,
                            "orientation": "horizontal", "min_width": 1024,
                            "safesearch": "true", "image_type": img_type},
                    timeout=15
                ) as r:
                    if r.status == 200:
                        hits = (await r.json()).get("hits", [])
                        if hits:
                            url = random.choice(hits)["largeImageURL"]
                            img = await _download(url)
                            if img:
                                log.info(f"✅ Pixabay: {query[:60]}")
                                return img
    except: pass
    return None

async def search_google(query: str) -> Optional[BytesIO]:
    if not GOOGLE_API_KEY or not GOOGLE_CX or not query: return None
    try:
        async with aiohttp.ClientSession() as s:
            async with s.get(
                "https://www.googleapis.com/customsearch/v1",
                params={"key": GOOGLE_API_KEY, "cx": GOOGLE_CX, "q": query,
                       "searchType": "image", "num": 5, "imgSize": "large", "safe": "active"},
                timeout=15
            ) as r:
                if r.status == 200:
                    items = (await r.json()).get("items", [])
                    if items:
                        url = random.choice(items)["link"]
                        img = await _download(url)
                        if img:
                            log.info(f"✅ Google: {query[:60]}")
                            return img
    except: pass
    return None

async def search_bing(query: str) -> Optional[BytesIO]:
    try:
        async with aiohttp.ClientSession() as s:
            async with s.get(
                f"https://www.bing.com/images/search?q={query.replace(' ', '+')}&qft=+filterui:imagesize-wallpaper",
                headers={"User-Agent": "Mozilla/5.0"},
                timeout=15
            ) as r:
                if r.status == 200:
                    html = await r.text()
                    urls = re.findall(r'https?://[^"\']+\.(?:jpg|jpeg|png|webp)', html)[:10]
                    if urls:
                        url = random.choice(urls)
                        img = await _download(url)
                        if img:
                            log.info(f"✅ Bing: {query[:60]}")
                            return img
    except: pass
    return None

async def get_image(visual_prompt: str) -> Optional[BytesIO]:
    clean = re.sub(r'--.*', '', visual_prompt).strip()
    words = clean.split()
    queries = [' '.join(words[:6])]
    if len(words) > 3:
        queries.append(' '.join(words[:3]))
    
    for query in queries:
        for search_fn in [search_unsplash, search_pixabay, search_google, search_bing]:
            img = await search_fn(query)
            if img: return img
    
    safe = re.sub(r'[^a-zA-Z0-9\s]', '', queries[0]).strip()
    if safe:
        try:
            async with aiohttp.ClientSession() as s:
                async with s.get(
                    f"https://image.pollinations.ai/prompt/{safe.replace(' ', '%20')}",
                    params={"width": 1024, "height": 768, "nologo": "true", "seed": str(random.randint(1, 99999))},
                    timeout=25
                ) as r:
                    if r.status == 200:
                        data = await r.read()
                        if len(data) > 2000:
                            return BytesIO(data)
        except: pass
    return None

# ===== PPTX С ШАБЛОННЫМИ СЛАЙДАМИ =====
async def make_pptx(data):
    slides = data.get("slides", [])
    about_text = data.get("about_text", "")
    topic = data.get("topic", "Презентация")
    
    if not slides:
        return None

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    
    # Выбираем случайный стиль
    theme = random.choice(PPT_THEMES)
    log.info(f"Стиль: {theme['name']}")

    # Готовим картинки для ВСЕХ содержательных слайдов
    tasks = [get_image(s.get("visual_prompt", "")) for s in slides]
    images = await asyncio.gather(*tasks, return_exceptions=True)

    # === СЛАЙД 1: ТИТУЛЬНЫЙ ===
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой макет
    _set_bg(sl, theme["bg"])
    
    # Заголовок
    title_box = sl.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    tf.text = topic
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*hex_to_rgb(theme["title"]))
    p.alignment = PP_ALIGN.CENTER
    
    # Подзаголовок
    sub_box = sl.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    tf2 = sub_box.text_frame
    tf2.text = f"Москва, {datetime.now().year}"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(*hex_to_rgb(theme["accent"]))
    p2.alignment = PP_ALIGN.CENTER

    # === СЛАЙД 2: "О ЧЁМ ЭТА ПРЕЗЕНТАЦИЯ?" ===
    sl2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl2, theme["bg"])
    
    q_box = sl2.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1.5))
    q_tf = q_box.text_frame
    q_tf.text = "О чём эта презентация?"
    q_p = q_tf.paragraphs[0]
    q_p.font.size = Pt(36)
    q_p.font.bold = True
    q_p.font.color.rgb = RGBColor(*hex_to_rgb(theme["accent"]))
    q_p.alignment = PP_ALIGN.CENTER
    
    a_box = sl2.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(3))
    a_tf = a_box.text_frame
    a_tf.word_wrap = True
    a_tf.text = about_text or "Ключевые аспекты и идеи по теме презентации."
    for p in a_tf.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(*hex_to_rgb(theme["text"]))
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)

    # === СОДЕРЖАТЕЛЬНЫЕ СЛАЙДЫ ===
    for i, s in enumerate(slides):
        img_on_left = (i % 2 == 0)
        
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(sl, theme["bg"])
        
        # Заголовок
        ttl = sl.shapes.add_textbox(MG, Emu(457200), SW - MG*2, Emu(914400))
        ttl_tf = ttl.text_frame
        ttl_tf.text = s.get("title", "Слайд")
        ttl_p = ttl_tf.paragraphs[0]
        ttl_p.font.size = Pt(30)
        ttl_p.font.bold = True
        ttl_p.font.color.rgb = RGBColor(*hex_to_rgb(theme["title"]))
        ttl_p.alignment = PP_ALIGN.LEFT
        
        # Текст
        if img_on_left:
            txt_left = MG + IW + GP
            img_left = MG
        else:
            txt_left = MG
            img_left = SW - IW - MG
        
        txBox = sl.shapes.add_textbox(txt_left, Emu(1600000), TW, Emu(4500000))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = s.get("expert_text", "")
        for p in tf.paragraphs:
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(*hex_to_rgb(theme["text"]))
            p.space_after = Pt(8)
        
        # Картинка
        img = images[i]
        if isinstance(img, BytesIO):
            try:
                sl.shapes.add_picture(img, img_left, Emu(1800000), width=IW)
            except:
                pass

    # === ПОСЛЕДНИЙ СЛАЙД: "СПАСИБО ЗА ВНИМАНИЕ!" ===
    sl_end = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl_end, theme["bg"])
    
    thanks_box = sl_end.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(3))
    thanks_tf = thanks_box.text_frame
    thanks_tf.text = "Спасибо за внимание!"
    thanks_p = thanks_tf.paragraphs[0]
    thanks_p.font.size = Pt(48)
    thanks_p.font.bold = True
    thanks_p.font.color.rgb = RGBColor(*hex_to_rgb(theme["accent"]))
    thanks_p.alignment = PP_ALIGN.CENTER
    
    # Подпись
    sub_end = sl_end.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.8))
    sub_tf = sub_end.text_frame
    sub_tf.text = f"© {datetime.now().year} | PrezaBot"
    sub_p = sub_tf.paragraphs[0]
    sub_p.font.size = Pt(16)
    sub_p.font.color.rgb = RGBColor(*hex_to_rgb(theme["text"]))
    sub_p.alignment = PP_ALIGN.CENTER

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

def _set_bg(slide, color_hex):
    """Устанавливает цвет фона слайда."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(color_hex))

def filename(topic):
    name = re.sub(r'[^\w\s-]', '', topic).strip()
    if len(name) > 30:
        name = name[:30].rsplit(' ', 1)[0]
    return f"{name.replace(' ', '_') or 'presentation'}.pptx"

async def send_file(msg, data, name, caption):
    for t in range(3):
        try:
            return await msg.answer_document(
                BufferedInputFile(data, name), caption=caption, parse_mode="Markdown"
            )
        except TelegramRetryAfter as e:
            await asyncio.sleep(e.retry_after)
        except:
            await asyncio.sleep(2)

# ========== ОБРАБОТЧИКИ ==========

@dp.message(Command("start"))
async def start(msg: Message, state: FSMContext):
    await state.clear()
    admin = "🆓 Бесплатный доступ!\n" if msg.from_user.id == ADMIN_ID else ""
    await msg.answer(
        f"🎓 *PrezaBot — презентации с ИИ!*\n\n"
        f"✨ Экспертный контент\n🖼️ Концептуальные фото\n🎨 Стильный дизайн\n\n"
        f"💰 {PRICE}₽\n{admin}👇",
        parse_mode="Markdown", reply_markup=menu()
    )

@dp.message(F.text == "ℹ️ Помощь")
async def help_cmd(msg: Message):
    await msg.answer(
        "📌 *Как создать:*\n\n"
        "1. Нажми «Создать»\n"
        "2. Напиши тему\n"
        "   `Нейросети` — бот сам решит сколько слайдов\n"
        "   `История 7` — 7 слайдов\n"
        "3. Оплати 100₽\n"
        "4. Получи файл!\n\n"
        "🎨 Слайды: тема → суть → контент → спасибо",
        parse_mode="Markdown"
    )

@dp.message(F.text == "💰 Цена")
async def price_cmd(msg: Message):
    await msg.answer(f"💎 *{PRICE}₽*\n✅ Экспертный текст\n✅ Концептуальные фото\n✅ Стильный дизайн\n✅ Шаблонные слайды", parse_mode="Markdown")

@dp.message(F.text == "🎨 Создать презентацию")
async def start_create(msg: Message, state: FSMContext):
    await state.clear()
    await state.set_state(State.topic)
    await msg.answer("✏️ *Напишите тему:*\n\nС числом: `Нейросети 7`\nБез числа: `История России`", parse_mode="Markdown")

@dp.message(StateFilter(State.topic))
async def got_topic(msg: Message, state: FSMContext):
    text = msg.text.strip()
    if text.startswith('/'):
        await state.clear()
        return await start(msg, state)
    
    parts = text.split()
    user_n = None
    try:
        n = int(parts[-1])
        if 3 <= n <= 12:
            user_n = n
            topic = " ".join(parts[:-1])
        else:
            topic = text
    except ValueError:
        topic = text
    
    await state.update_data(topic=topic, num=user_n)
    
    if msg.from_user.id == ADMIN_ID:
        await state.clear()
        n_text = f"{user_n} слайдов" if user_n else "оптимально"
        status = await msg.answer(f"🔄 Генерирую «{topic}», {n_text}...")
        try:
            data = await asyncio.wait_for(get_content(topic, user_n), timeout=180)
            if not data:
                return await status.edit_text("❌ ИИ не отвечает.")
            total = len(data.get("slides", []))
            await status.edit_text(f"🔍 {total} слайдов. Подбираю фото...")
            pptx = await asyncio.wait_for(make_pptx(data), timeout=300)
            if not pptx:
                return await status.edit_text("❌ Ошибка сборки")
            await send_file(msg, pptx.getvalue(), filename(topic),
                          f"✅ *Готово!*\n📌 {topic}\n📊 {total} слайдов + шаблонные")
            await status.delete()
        except asyncio.TimeoutError:
            await status.edit_text("⏰ Долго.")
        except Exception as e:
            log.error(f"{e}")
            await status.edit_text("❌ Ошибка. /start")
    else:
        try:
            payment = Payment.create({
                "amount": {"value": f"{PRICE}.00", "currency": "RUB"},
                "confirmation": {"type": "redirect", "return_url": f"https://t.me/{(await bot.get_me()).username}"},
                "description": f"Презентация «{topic[:50]}»",
                "metadata": {"uid": msg.from_user.id, "topic": topic, "n": user_n},
                "capture": True,
                "receipt": {
                    "customer": {"email": f"{msg.from_user.id}@t.me"},
                    "items": [{"description": f"Презентация «{topic[:30]}»", "quantity": "1",
                              "amount": {"value": f"{PRICE}.00", "currency": "RUB"},
                              "vat_code": "1", "payment_mode": "full_prepayment", "payment_subject": "service"}]
                }
            })
            await state.update_data(pid=payment.id)
            await state.set_state(State.payment)
            await msg.answer(f"💎 *Заказ*\n📌 {topic}\n💰 *{PRICE}₽*\n👇 Оплатить:",
                           parse_mode="Markdown", reply_markup=pay_kb(payment.confirmation.confirmation_url))
        except Exception as e:
            log.error(f"Платёж: {e}")
            await msg.answer("❌ Ошибка")
            await state.clear()

@dp.callback_query(F.data == "paid")
async def check_pay(cb: CallbackQuery, state: FSMContext):
    d = await state.get_data()
    pid = d.get("pid")
    if not pid: return await cb.answer("❌ Нет платежа")
    try:
        p = Payment.find_one(pid)
    except:
        return await cb.answer("❌ Ошибка")
    if p.status == "succeeded":
        topic = d.get("topic") or (p.metadata or {}).get("topic")
        user_n = d.get("num") or (p.metadata or {}).get("n")
        await cb.message.edit_text(f"✅ Оплачено! Генерирую «{topic}»...")
        try:
            data = await asyncio.wait_for(get_content(topic, user_n), timeout=180)
            if not data: return await cb.message.edit_text("❌ ИИ не ответил")
            total = len(data.get("slides", []))
            await cb.message.edit_text(f"🔍 {total} слайдов. Подбираю фото...")
            pptx = await asyncio.wait_for(make_pptx(data), timeout=300)
            if not pptx: return await cb.message.edit_text("❌ Ошибка сборки")
            await send_file(cb.message, pptx.getvalue(), filename(topic),
                          f"✅ *Готово!*\n📌 {topic}\n📊 {total} слайдов + шаблонные\n💰 {PRICE}₽ оплачено")
            await cb.message.delete()
        except asyncio.TimeoutError:
            await cb.message.edit_text("⏰ Долго")
        except Exception as e:
            log.error(f"{e}")
            await cb.message.edit_text("❌ Ошибка")
        await state.clear()
    elif p.status == "pending":
        await cb.answer("⏳ Жди", show_alert=True)
    else:
        await cb.answer(f"❌ {p.status}", show_alert=True)
        await state.clear()

@dp.callback_query(F.data == "cancel")
async def cancel_pay(cb: CallbackQuery, state: FSMContext):
    await state.clear()
    await cb.message.edit_text("❌ Отменено. /start")

@dp.message()
async def fallback(msg: Message):
    await msg.answer("🤔 Меню или /start", reply_markup=menu())

async def main():
    log.info("🚀 Запуск с шаблонными слайдами и стилями")
    await bot.delete_webhook(drop_pending_updates=True)
    await dp.start_polling(bot, allowed_updates=["message", "callback_query"])

if __name__ == "__main__":
    asyncio.run(main())